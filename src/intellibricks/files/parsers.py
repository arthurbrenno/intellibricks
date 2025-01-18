from __future__ import annotations

import abc
import io
import tempfile
from typing import Optional, Sequence, TypedDict, cast, override

import msgspec
from architecture.data.files import FileExtension, RawFile
from architecture.utils.decorators import ensure_module_installed
from architecture.utils.functions import run_sync
from openai import OpenAI

from intellibricks.agents import Agent
from intellibricks.llms.types import (
    ChainOfThought,
    ImageDescription,
    ImageFilePart,
    MimeType,
)

from .constants import ParsingStrategy
from .parsed_files import Image, PageContent, ParsedFile


class LocalSettings(TypedDict):
    use_gpu: bool


class FileParser(msgspec.Struct, frozen=True):
    """
    Abstract class for extracting content from files.
    This should be used as a base class for specific file parsers.
    """

    strategy: ParsingStrategy = ParsingStrategy.DEFAULT

    def extract_contents(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        return run_sync(self.extract_contents_async, file)

    @abc.abstractmethod
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        raise NotImplementedError("This method should be implemented by subclasses.")


class IntellibricksFileParser(FileParser, frozen=True):
    image_description_agent: Optional[Agent[ChainOfThought[ImageDescription]]] = None


class PDFFileParser(IntellibricksFileParser, frozen=True):
    @ensure_module_installed("pypdf", "intellibricks[files]")
    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        from pypdf import PdfReader

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            reader = PdfReader(file_path)
            page_contents: list[PageContent] = []
            for page_num, page in enumerate(reader.pages):
                page_images: Sequence[Image] = [
                    Image(contents=image.data, name=image.name) for image in page.images
                ]

                image_descriptions: list[str] = []
                if (
                    self.image_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    for image_num, image in enumerate(page_images):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png, data=image.contents
                        )
                        agent_response = await self.image_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Page Image {image_num + 1}: {image_md}"
                        )

                page_text = [page.extract_text(), "".join(image_descriptions)]

                page_content = PageContent(
                    page=page_num + 1,
                    text="".join(page_text),
                    images=page_images,
                )

                page_contents.append(page_content)

            file_name = file.name
            return ParsedFile(
                name=file_name,
                pages=page_contents,
            )


class OfficeFileParser(IntellibricksFileParser, frozen=True):
    """
    This class actually delegates the parsing to the appropriate parser based on the file extension.
    This class is a Facade for the different Office file parsers.
    """

    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        extension = file.extension
        match extension:
            case FileExtension.DOCX:
                return await DocxFileParser(
                    strategy=self.strategy
                ).extract_contents_async(file)
            case FileExtension.PPTX:
                return await PptxFileParser(
                    strategy=self.strategy
                ).extract_contents_async(file)
            case FileExtension.XLSX:
                return await ExcelFileParser(
                    strategy=self.strategy
                ).extract_contents_async(file)
            case _:
                raise ValueError(f"Unsupported file extension: {extension}")


class DocxFileParser(OfficeFileParser, frozen=True):
    @ensure_module_installed("python-docx", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import tempfile

        from docx import Document  # python-docx

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            document = Document(file_path)

            # Extract all text from paragraphs
            paragraph_texts: list[str] = []
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    paragraph_texts.append(paragraph.text)
            doc_text = "\n".join(paragraph_texts)

            # Extract all images
            doc_images: list[Image] = []
            for rel in document.part._rels.values():  # type: ignore
                # Relationship is image-based if it references an image part
                if "image" in rel.reltype:
                    image_part = rel.target_part
                    image_name = image_part.partname.split("/")[-1]  # e.g. "image1.png"
                    image_bytes = image_part.blob
                    doc_images.append(Image(name=image_name, contents=image_bytes))

            # If high-level strategy, describe images
            image_descriptions: list[str] = []
            if self.image_description_agent and self.strategy == ParsingStrategy.HIGH:
                for idx, image in enumerate(doc_images, start=1):
                    agent_input = ImageFilePart(
                        mime_type=MimeType.image_png,  # or detect from extension
                        data=image.contents,
                    )
                    agent_response = await self.image_description_agent.run_async(
                        agent_input
                    )
                    image_md = agent_response.parsed.final_answer.md
                    image_descriptions.append(f"Docx Image {idx}: {image_md}")

                # Append the images' descriptions to the main text
                if image_descriptions:
                    doc_text += "\n\n" + "\n".join(image_descriptions)

            # Create a single PageContent (DOCX has no true "pages" by default)
            page_content = PageContent(
                page=1,
                text=doc_text,
                images=doc_images,
            )

            return ParsedFile(
                name=file.name,
                pages=[page_content],
            )


class PptxFileParser(OfficeFileParser, frozen=True):
    @ensure_module_installed("python-pptx", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import tempfile

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            prs: PptxPresentation = Presentation(file_path)

            pages: list[PageContent] = []

            for slide_index, slide in enumerate(prs.slides, start=1):
                # We'll store text from shapes and images
                slide_texts: list[str] = []
                slide_images: list[Image] = []

                # Examine each shape
                for shape in slide.shapes:
                    # If shape has a text frame, cast to Shape
                    if shape.has_text_frame:
                        shape_with_text = cast(Shape, shape)
                        text_str: str = shape_with_text.text
                        slide_texts.append(text_str)

                    # If shape is a picture, cast to Picture
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_shape = cast(Picture, shape)
                        image_blob: bytes = picture_shape.image.blob
                        image_name: str = shape.name or f"slide_{slide_index}_img"
                        slide_images.append(Image(name=image_name, contents=image_blob))

                combined_text: str = "\n".join(slide_texts)

                # If strategy is HIGH, we generate image descriptions
                if (
                    self.image_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(slide_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png,
                            data=image_obj.contents,
                        )
                        agent_response = await self.image_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {image_md}"
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                page_content = PageContent(
                    page=slide_index,
                    text=combined_text,
                    images=slide_images,
                )
                pages.append(page_content)

            return ParsedFile(
                name=file.name,
                pages=pages,
            )


class ExcelFileParser(OfficeFileParser, frozen=True):
    @ensure_module_installed("openpyxl", "intellibricks[files]")
    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import tempfile

        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            wb: Workbook = load_workbook(file_path, data_only=True)
            pages: list[PageContent] = []

            for sheet_index, sheet in enumerate(wb.worksheets, start=1):
                # Gather all text
                row_texts: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    # row might be a tuple of cell values, some may be None
                    cell_strs = [str(cell) for cell in row if cell is not None]
                    if cell_strs:
                        row_texts.append("\t".join(cell_strs))

                combined_text: str = "\n".join(row_texts)

                # Gather images (via private _images)
                sheet_images: list[Image] = []
                # This is a private attribute in openpyxl, so type stubs are nonexistent.
                # We'll do a "if hasattr(sheet, '_images')" to avoid exceptions.
                if hasattr(sheet, "_images"):
                    image_list = getattr(sheet, "_images", [])
                    for img_idx, img in enumerate(image_list, start=1):
                        # openpyxl image objects often have a "ref" attribute, plus a "_data" attribute for bytes
                        # or ._data is sometimes called .image depending on your openpyxl version
                        # We can do a cast or just ignore
                        img_data = getattr(img, "_data", None)
                        if img_data is not None:
                            image_name = f"{sheet.title}_img_{img_idx}.png"
                            sheet_images.append(
                                Image(name=image_name, contents=img_data)
                            )

                # If strategy is HIGH, run image_description_agent
                if (
                    self.image_description_agent
                    and self.strategy == ParsingStrategy.HIGH
                ):
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(sheet_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=MimeType.image_png,
                            data=image_obj.contents,
                        )
                        agent_response = await self.image_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Worksheet {sheet.title} - Image {img_idx}: {image_md}"
                        )
                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                page_content = PageContent(
                    page=sheet_index,
                    text=combined_text,
                    images=sheet_images,
                )
                pages.append(page_content)

            return ParsedFile(
                name=file.name,
                pages=pages,
            )


class TxtFileParser(IntellibricksFileParser, frozen=True):
    """
    Parses plain .txt files. Extracts all content as a single page (page=1).
    """

    @override
    async def extract_contents_async(self, file: RawFile) -> ParsedFile:
        text_content = file.contents.decode("utf-8")

        page_content = PageContent(
            page=1,
            text=text_content,
            md=text_content,
        )

        return ParsedFile(
            name=file.name,
            pages=[page_content],
        )


class StaticImageFileParser(IntellibricksFileParser, frozen=True):
    """
    Parses static image files (PNG, JPEG, TIFF, etc.) as a single "page" with one image.
    If the image is TIFF, it converts to PNG in-memory for better compatibility.
    If the strategy == HIGH and an image_description_agent is present,
    it appends an AI-generated textual description of the image.
    """

    @override
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Determine the extension
            extension = file.extension.value.lower()  # e.g. "png", "jpg", "tiff"

            # Convert to PNG if TIFF
            if extension in {"tiff", "tif"}:
                # Use Pillow to open, then convert to PNG in memory
                with io.BytesIO(file.contents) as input_buffer:
                    with PILImage.open(input_buffer) as pil_img:
                        # Convert to RGBA or RGB if needed
                        if pil_img.mode not in ("RGB", "RGBA"):
                            pil_img = pil_img.convert("RGBA")

                        # Save as PNG into a new buffer
                        output_buffer = io.BytesIO()
                        pil_img.save(output_buffer, format="PNG")
                        converted_bytes = output_buffer.getvalue()

                # Use the converted PNG bytes
                image_bytes = converted_bytes
                current_mime_type = MimeType.image_png
            else:
                # No conversion needed
                image_bytes = file.contents

                # For demonstration, pick your MIME by extension
                if extension in {"png"}:
                    current_mime_type = MimeType.image_png
                elif extension in {"jpg", "jpeg"}:
                    current_mime_type = MimeType.image_jpeg
                else:
                    # Fallback to PNG or raise an error if you want
                    current_mime_type = MimeType.image_png

            # Create an Image object
            image_obj = Image(name=file.name, contents=image_bytes)

            # Generate a description if we have an agent + HIGH strategy
            text_content = ""
            if self.image_description_agent and self.strategy == ParsingStrategy.HIGH:
                agent_input = ImageFilePart(
                    mime_type=current_mime_type,
                    data=image_bytes,
                )
                agent_response = await self.image_description_agent.run_async(
                    agent_input
                )
                description_md = agent_response.parsed.final_answer.md
                text_content = description_md

            # We treat it as a single "page" with one image
            page_content = PageContent(
                page=1,
                text=text_content,
                images=[image_obj],
            )

            return ParsedFile(
                name=file.name,
                pages=[page_content],
            )


# ['doc', 'docx', 'txt', 'pdf', 'xlsx', 'xls', 'jpg', 'jpeg', 'tif', 'tiff', 'bmp', 'png', 'PNG', 'gif', 'ppt', 'pptx', 'pptm', 'pkt', 'alg', 'pkz', 'rar', 'zip', 'dwg']


class MarkitdownFileParser(FileParser, frozen=True):
    client: Optional[OpenAI] = None
    model: Optional[str] = None

    @ensure_module_installed("markitdown", "intellibricks[files]")
    async def extract_contents_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from markitdown import MarkItDown
        from markitdown._markitdown import DocumentConverterResult

        match self.strategy:
            case (
                ParsingStrategy.DEFAULT | ParsingStrategy.MEDIUM | ParsingStrategy.FAST
            ):
                llm_client = None
                llm_model = None
            case ParsingStrategy.HIGH:
                llm_client = self.client or OpenAI()
                llm_model = self.model or "gpt-4o"

        with tempfile.NamedTemporaryFile(delete=True) as temp_file:
            temp_file.write(file.contents)
            temp_file.seek(0)
            converter = MarkItDown(llm_client=llm_client, llm_model=llm_model)
            result: DocumentConverterResult = converter.convert(temp_file.name)
            markdown: str = result.text_content

            # return a Document with one page only
            page_content = PageContent(
                page=1,
                text=markdown,
                md=markdown,
                images=[],
                items=[],
            )

            return ParsedFile(
                name=file.name,
                pages=[page_content],
            )
